import json
from pathlib import Path

import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


def parse_uploaded_file(file_path: Path, original_name: str, content_type: str):
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        text = extract_pdf_text(file_path)
    elif suffix == ".docx":
        text = extract_docx_text(file_path)
    elif suffix in [".pptx", ".ppt"]:
        text = extract_ppt_text(file_path) if suffix == ".pptx" else ""
    elif suffix == ".txt":
        text = file_path.read_text(encoding="utf-8", errors="ignore")
    elif suffix == ".csv":
        text = extract_csv_text(file_path)
    elif suffix in [".xlsx", ".xls"]:
        text = extract_excel_text(file_path)
    elif suffix == ".json":
        text = extract_json_text(file_path)
    elif suffix == ".doc":
        # Old .doc needs external conversion tools. Convert to .docx for best result.
        text = ""
    else:
        text = ""

    return {
        "source_type": "file",
        "content_type": content_type,
        "file_name": original_name,
        "url": None,
        "title": original_name,
        "text": text,
    }


def extract_pdf_text(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    pages = []

    for page in reader.pages:
        pages.append(page.extract_text() or "")

    return "\n".join(pages)


def extract_docx_text(file_path: Path) -> str:
    doc = Document(str(file_path))
    lines = []

    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                lines.append(row_text)

    return "\n".join(lines)


def extract_csv_text(file_path: Path) -> str:
    df = pd.read_csv(file_path)
    return df.astype(str).to_csv(index=False)


def extract_excel_text(file_path: Path) -> str:
    excel = pd.ExcelFile(file_path)
    parts = []

    for sheet_name in excel.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        parts.append(f"Sheet: {sheet_name}\n")
        parts.append(df.astype(str).to_csv(index=False))

    return "\n\n".join(parts)


def extract_json_text(file_path: Path) -> str:
    data = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(data, ensure_ascii=False, indent=2)


def extract_ppt_text(file_path: Path) -> str:
    prs = Presentation(str(file_path))
    lines = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_lines.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_lines.append(row_text)
        if slide_lines:
            lines.append(f"Slide {slide_no}:\n" + "\n".join(slide_lines))
    return "\n\n".join(lines)
